# agents/presentation_agent.py
# PresentationAgent updated to add speaker notes to slides.

from .base_agent import BaseAgent
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE # Correct Enum import
import os
import streamlit as st

class PresentationAgent(BaseAgent):
    """
    Generates the final .pptx presentation, including speaker notes,
    handling layouts, images, and attempting to auto-fit text.
    """

    def _delete_initial_slide(self, prs, slides_plan):
        # ... (delete_initial_slide remains the same)
        while len(prs.slides) > len(slides_plan):
            xml_slides = prs.slides._sldIdLst
            to_remove = xml_slides[0]; xml_slides.remove(to_remove)
            self.log("Removed an initial blank slide.")

    def _set_font_size(self, text_frame, size_pt):
        # ... (set_font_size remains the same)
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size_pt)

    def run(self):
        self.log("Starting final presentation generation...")
        slides_plan = self.sm.get("slides") or []
        design_config = self.sm.get("design")
        output_dir = "output"
        output_filename = "final_presentation.pptx"
        output_path = os.path.join(output_dir, output_filename)

        if not slides_plan:
            self.log("ERROR: No slides plan found."); return

        template_path = design_config.get("template_path")
        try:
            prs = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()
            self.log(f"Using template: {template_path}" if template_path and os.path.exists(template_path) else "Default presentation.")
        except Exception as e:
            self.log(f"ERROR: Failed loading template '{template_path}'. {e}"); prs = Presentation()

        layout_map = { "main_title": 0, "chapter_title": 0, "content_only": 1, "content_with_image": 8, "quiz": 1, "thank_you": 5 }

        for slide_data in slides_plan:
            slide_type = slide_data.get("type", "content")
            image_path = slide_data.get("image_path")
            layout_key = "content_only"
            if slide_type == "content" and image_path and os.path.exists(image_path): layout_key = "content_with_image"
            elif slide_type != "content": layout_key = slide_type

            layout_index = layout_map.get(layout_key, 1)

            try:
                slide_layout = prs.slide_layouts[layout_index]
                slide = prs.slides.add_slide(slide_layout)
            except IndexError:
                self.log(f"WARN: Layout index {layout_index} not found. Using layout 1.");
                slide_layout = prs.slide_layouts[1]; slide = prs.slides.add_slide(slide_layout)

            # --- ADD SPEAKER NOTES ---
            # Check if the slide object supports notes (most layouts do)
            # and if speaker notes data exists in our slide_data dictionary.
            if slide.has_notes_slide and slide_data.get("speaker_notes"):
                notes_slide = slide.notes_slide # Get the notes slide object associated with this main slide.
                text_frame = notes_slide.notes_text_frame # Get the text frame within the notes slide.
                text_frame.clear() # Clear any default placeholder text in the notes.
                p = text_frame.add_paragraph() # Add a new paragraph to the notes.
                p.text = slide_data.get("speaker_notes", "") # Set the paragraph text to the notes from slide_data.
                self.log(f"Added speaker notes to slide '{slide_data.get('title', 'Untitled')}'.") # Log confirmation.
            # -------------------------

            # Populate Title (remains the same)
            if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
                title_shape = slide.shapes.title; title_shape.text = slide_data.get("title", "")
                title_shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                title_shape.text_frame.word_wrap = True

            # Populate Content Placeholders (remains the same)
            if layout_key in ["main_title", "chapter_title"]:
                if len(slide.placeholders) > 1:
                    sub_ph = slide.placeholders[1]; sub_ph.text = slide_data.get("subtitle", "")
                    sub_ph.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    sub_ph.text_frame.word_wrap = True

            elif layout_key in ["content_only", "quiz"]:
                if len(slide.placeholders) > 1:
                    body_ph = slide.placeholders[1]; tf = body_ph.text_frame; tf.clear()
                    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    self._set_font_size(tf, 16)
                    for bullet in slide_data.get("bullets", []):
                        p = tf.add_paragraph()
                        if isinstance(bullet, dict): p.text = bullet.get('question', '')
                        else: p.text = str(bullet)
                        p.level = 0
                        for run in p.runs: run.font.size = Pt(16)

            elif layout_key == "content_with_image":
                if len(slide.placeholders) > 2:
                    text_ph = slide.placeholders[1]; tf = text_ph.text_frame; tf.clear()
                    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    self._set_font_size(tf, 14)
                    for bullet in slide_data.get("bullets", []):
                        p = tf.add_paragraph(); p.text = str(bullet); p.level = 0
                        for run in p.runs: run.font.size = Pt(14)

                    img_ph = slide.placeholders[2]
                    if image_path and os.path.exists(image_path):
                        try:
                            slide.shapes.add_picture(
                                image_path, img_ph.left, img_ph.top,
                                width=img_ph.width, height=img_ph.height)
                        except Exception as img_e: self.log(f"ERROR adding picture {image_path}: {img_e}")

        self._delete_initial_slide(prs, slides_plan)
        os.makedirs(output_dir, exist_ok=True)
        try:
            prs.save(output_path)
            self.update_state("output_path", output_path)
            self.log(f"Presentation saved successfully: {output_path}")
        except PermissionError:
             self.log(f"ERROR: Permission denied saving {output_path}. Is file open?")
             st.error(f"Save failed: Close {os.path.basename(output_path)} and retry.")
        except Exception as save_e:
             self.log(f"ERROR saving presentation: {save_e}")
             st.error(f"Failed to save presentation: {save_e}")
























































































































# from state_manager import StateManager
# from pptx import Presentation
# from pptx.util import Inches, Pt
# import logging
# import os

# class PresentationAgent:
#     def __init__(self, state: StateManager, output_path="output_ppt/final_presentation.pptx"):
#         self.state = state
#         self.output_path = output_path
#         logging.basicConfig(level=logging.INFO, format='[%(asctime)s] [%(name)s] %(message)s')
#         self.logger = logging.getLogger("PresentationAgent")

#     def build_presentation(self):
#         self.logger.info("Starting presentation generation...")

#         # ✅ Correct way — StateManager.get() only takes one argument
#         slides = self.state.get("slides")
#         design = self.state.get("design")

#         # Provide safe defaults manually if None
#         if slides is None:
#             slides = []
#         if design is None:
#             design = {}

#         if not slides:
#             self.logger.warning("No slides found in shared state. Aborting presentation build.")
#             return

#         prs = Presentation()

#         for idx, slide_content in enumerate(slides):
#             slide_layout = prs.slide_layouts[1]  # Title + content layout
#             slide = prs.slides.add_slide(slide_layout)

#             title = slide.shapes.title
#             content = slide.placeholders[1]

#             title.text = slide_content.get("title", f"Slide {idx + 1}")
#             content.text = slide_content.get("content", "No content provided.")

#             # Optionally apply dummy design style
#             style = design.get(f"slide_{idx+1}", "Default Style")
#             self.logger.info(f"Applied design style '{style}' to slide {idx + 1}")

#         # Ensure output directory exists
#         os.makedirs(os.path.dirname(self.output_path), exist_ok=True)

#         prs.save(self.output_path)
#         self.logger.info(f"Presentation saved at: {self.output_path}")

#_________________________________________________________________________________________
